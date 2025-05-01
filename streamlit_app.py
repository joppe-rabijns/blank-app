import streamlit as st
import pandas as pd
from pptx import Presentation
import tempfile
import pycountry

st.title("Punten 2025 PowerPoint Generator")

uploaded_file = st.file_uploader("Upload the Excel file", type=["xlsx"])
template_file = st.file_uploader("Upload a PPTX template", type=["pptx"])

# Convert 2-letter country code to full name
def get_country_name(code):
    try:
        return pycountry.countries.get(alpha_2=code.upper()).name
    except:
        return code

# Translate score categories
def translate_prize(code):
    if pd.isna(code):
        return ""
    code = str(code).strip()
    return {
        "1 CL": "FIRST PRIZE\nCUM LAUDE",
        "1 SCL": "FIRST PRIZE\nSUMMA CUM LAUDE",
        "1": "FIRST PRIZE",
        "2": "SECOND PRIZE",
        "3": "THIRD PRIZE",
        "Certificate of participation": "CERTIFICATE OF PARTICIPATION",
        "mentioned": "MENTIONED"
    }.get(code, code)

# Helper to fill placeholders by index
def fill_placeholders_by_index(slide, mapping: dict):
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if not shape.has_text_frame:
            continue
        if idx in mapping:
            shape.text = str(mapping[idx])

if uploaded_file and template_file:
    xls = pd.ExcelFile(uploaded_file)
    sheet_names = xls.sheet_names
    selected_sheet = st.selectbox("Select the sheet to generate slides from", sheet_names)

    df_preview = pd.read_excel(uploaded_file, sheet_name=selected_sheet)
    st.subheader(f"Preview of {selected_sheet}")
    st.dataframe(df_preview)

    prs = Presentation(template_file)

    # Remove the first slide from the template if present
    #if prs.slides:
    #    xml_slides = prs.slides._sldIdLst
    #    slides = list(xml_slides)
    #    prs.slides._sldIdLst.remove(slides[0])

    header_layout = prs.slide_layouts[0]
    participant_layout = prs.slide_layouts[1]

    df = pd.read_excel(uploaded_file, sheet_name=selected_sheet)
    day = selected_sheet.replace("Punten ", "")
    df["_order"] = range(len(df))
    grouped = df.groupby(["Lokatie", "Reeks"], sort=False)

    for (location, reeks), group in grouped:
        group = group.sort_values("_order")

        # Header slide
        header_slide = prs.slides.add_slide(header_layout)
        fill_placeholders_by_index(header_slide, {
            21: f"{location}",
            1: f"{reeks}"
        })

        for _, row in group.iterrows():
            naam = row.get("Naam", "")
            stad = row.get("Stad", "")
            land_code = str(row.get("Land", "")).strip()
            land = get_country_name(land_code)
            prijs_code = row.get("Prijscategorie", "")
            prijs = translate_prize(prijs_code)

            # Slide without prize
            slide1 = prs.slides.add_slide(participant_layout)

            formatted_reeks = reeks

            if ":" in reeks:
                parts = reeks.split(":", 1)
                formatted_reeks = f"{parts[0].strip()}\n{parts[1].strip()}"

            fill_placeholders_by_index(slide1, {
                 1: f"",
                 21: stad,
                 22: land,
                 23: formatted_reeks,
                 24: naam,
            })

            # Slide with prize
            slide2 = prs.slides.add_slide(participant_layout)
            fill_placeholders_by_index(slide2, {
                1: prijs,
                21: stad,
                22: land,
                23: formatted_reeks,
                24: naam,
            })

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        st.success("PowerPoint presentation generated successfully!")
        with open(tmp.name, "rb") as file:
            st.download_button(
                label="Download Presentation",
                data=file,
                file_name="Punten_Presentatie_2025.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )